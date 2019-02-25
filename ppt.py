def ppt():
    import os
    from pptx import Presentation 
    from pptx.util import Inches

    # initiate files
    files = os.listdir()
    pngfiles = [i for i in files if i.endswith('.png')]

    #initiate presentation
    prs = Presentation() 
    blank_slide_layout = prs.slide_layouts[6] 

    for i in pngfiles:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(1)
        pic = slide.shapes.add_picture(i, left, top)

    #saveppt
    prs.save('Presentation.pptx')